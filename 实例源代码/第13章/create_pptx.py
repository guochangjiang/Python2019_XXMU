# create_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
# 标题页
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Python制作精美PPT"
subtitle.text = "python-pptx"

# 列表页
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = '列表页'
tf = body_shape.text_frame
tf.text = '列表级别0'

p = tf.add_paragraph()
p.text = '列表级别1'
p.level = 1
p = tf.add_paragraph()
p.text = '列表级别2'
p.level = 2
p = tf.add_paragraph()
p.text = '列表级别2'
p.level = 3

# 文本框页
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "普通文本框"
p = tf.add_paragraph()
p.text = "文本框文字加粗"
p.font.bold = True
p = tf.add_paragraph()
p.text = "文本框大号文字"
p.font.size = Pt(40)

# 图片页
img_path = 'Python.png'
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left = top = Inches(1)
width = height = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, width, height)
left, top = Inches(5), Inches(1)
width, height = Inches(4), Inches(2)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 添加形状
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = '图形页'
left = Inches(0.93)  # 0.93" centers this overall set of shapes
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)
shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'
left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance
for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = 'Step %d' % n
    left = left + width - Inches(0.4)

# 添加表格
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = '表格页'
rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)
table = shapes.add_table(rows, cols, left, top, width, height).table
# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)
# write column headings
table.cell(0, 0).text = '行1列1'
table.cell(0, 1).text = '行1列2'
# write body cells
table.cell(1, 0).text = '行2列1'
table.cell(1, 1).text = '行2列2'
prs.save('demo.pptx')